"""
Document processor for extracting text from various file types.
Supports PDF, DOCX, TXT, MD, and common code files.
"""

import re
from pathlib import Path
from typing import Tuple, Dict, List, Optional
from datetime import datetime


class DocumentProcessor:
    """Process and extract text from various document types."""
    
    SUPPORTED_EXTENSIONS = {
        # Documents
        ".txt", ".md", ".markdown", ".rst", ".rtf",
        # Office documents
        ".pdf", ".docx", ".doc", ".pptx", ".xlsx",
        # Code files
        ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cpp", ".c", ".h",
        ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
        ".html", ".css", ".scss", ".sass", ".less",
        ".json", ".yaml", ".yml", ".toml", ".xml", ".ini", ".cfg",
        ".sh", ".bash", ".zsh", ".fish", ".ps1",
        ".sql", ".graphql",
        # Data files
        ".csv", ".tsv",
        # Notebooks
        ".ipynb",
    }
    
    CODE_EXTENSIONS = {
        ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".cpp", ".c", ".h",
        ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala",
        ".html", ".css", ".scss", ".sass", ".less", ".sh", ".bash",
        ".sql", ".graphql", ".json", ".yaml", ".yml", ".toml", ".xml",
    }
    
    def is_supported(self, path: Path) -> bool:
        """Check if file type is supported."""
        return path.suffix.lower() in self.SUPPORTED_EXTENSIONS
    
    def process(self, path: Path) -> Tuple[str, Dict]:
        """
        Extract text content and metadata from a file.
        
        Returns:
            Tuple of (content_string, metadata_dict)
        """
        suffix = path.suffix.lower()
        
        metadata = {
            "file_size": path.stat().st_size,
            "modified_at": datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
            "created_at": datetime.fromtimestamp(path.stat().st_ctime).isoformat(),
        }
        
        try:
            if suffix == ".pdf":
                content = self._process_pdf(path)
            elif suffix == ".docx":
                content = self._process_docx(path)
            elif suffix == ".pptx":
                content = self._process_pptx(path)
            elif suffix == ".xlsx":
                content = self._process_xlsx(path)
            elif suffix == ".ipynb":
                content = self._process_notebook(path)
            elif suffix in self.CODE_EXTENSIONS:
                content = self._process_code(path)
                metadata["is_code"] = True
            else:
                content = self._process_text(path)
            
            # Clean up content
            content = self._clean_text(content)
            metadata["char_count"] = len(content)
            metadata["word_count"] = len(content.split())
            
            return content, metadata
            
        except Exception as e:
            return "", {"error": str(e)}
    
    def _process_text(self, path: Path) -> str:
        """Process plain text files."""
        try:
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return path.read_text(encoding="latin-1")
    
    def _process_code(self, path: Path) -> str:
        """Process code files with language context."""
        content = self._process_text(path)
        
        # Add file context for better embeddings
        lang_map = {
            ".py": "Python",
            ".js": "JavaScript",
            ".ts": "TypeScript",
            ".jsx": "React JSX",
            ".tsx": "React TSX",
            ".java": "Java",
            ".cpp": "C++",
            ".c": "C",
            ".go": "Go",
            ".rs": "Rust",
            ".rb": "Ruby",
            ".php": "PHP",
            ".swift": "Swift",
            ".kt": "Kotlin",
            ".sql": "SQL",
        }
        
        lang = lang_map.get(path.suffix.lower(), "code")
        return f"[{lang} file: {path.name}]\n\n{content}"
    
    def _process_pdf(self, path: Path) -> str:
        """Process PDF files."""
        try:
            import pypdf
            
            reader = pypdf.PdfReader(str(path))
            text_parts = []
            
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    text_parts.append(f"[Page {page_num + 1}]\n{text}")
            
            return "\n\n".join(text_parts)
        except ImportError:
            # Fallback: try pdfplumber
            try:
                import pdfplumber
                
                with pdfplumber.open(str(path)) as pdf:
                    text_parts = []
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            text_parts.append(f"[Page {page_num + 1}]\n{text}")
                    return "\n\n".join(text_parts)
            except ImportError:
                return f"[PDF file: {path.name} - install pypdf or pdfplumber to extract text]"
    
    def _process_docx(self, path: Path) -> str:
        """Process Word documents."""
        try:
            from docx import Document
            
            doc = Document(str(path))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            return f"[DOCX file: {path.name} - install python-docx to extract text]"
    
    def _process_pptx(self, path: Path) -> str:
        """Process PowerPoint files."""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(path))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"[Slide {slide_num + 1}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts)
        except ImportError:
            return f"[PPTX file: {path.name} - install python-pptx to extract text]"
    
    def _process_xlsx(self, path: Path) -> str:
        """Process Excel files."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(str(path), data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        sheet_text.append(" | ".join(row_values))
                
                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))
            
            return "\n\n".join(text_parts)
        except ImportError:
            return f"[XLSX file: {path.name} - install openpyxl to extract text]"
    
    def _process_notebook(self, path: Path) -> str:
        """Process Jupyter notebooks."""
        import json
        
        try:
            notebook = json.loads(path.read_text())
            text_parts = []
            
            for cell_num, cell in enumerate(notebook.get("cells", [])):
                cell_type = cell.get("cell_type", "")
                source = "".join(cell.get("source", []))
                
                if cell_type == "markdown":
                    text_parts.append(f"[Markdown Cell {cell_num + 1}]\n{source}")
                elif cell_type == "code":
                    text_parts.append(f"[Code Cell {cell_num + 1}]\n```\n{source}\n```")
                    
                    # Include output if present
                    outputs = cell.get("outputs", [])
                    for output in outputs:
                        if "text" in output:
                            output_text = "".join(output["text"])
                            text_parts.append(f"Output:\n{output_text}")
            
            return "\n\n".join(text_parts)
        except Exception as e:
            return f"[Notebook file: {path.name} - error: {e}]"
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove excessive whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)
        text = re.sub(r'\t+', ' ', text)
        
        # Remove null bytes and other problematic characters
        text = text.replace('\x00', '')
        
        return text.strip()
    
    def chunk_text(
        self, 
        text: str, 
        chunk_size: int = 1000, 
        overlap: int = 200,
        respect_boundaries: bool = True
    ) -> List[str]:
        """
        Split text into overlapping chunks for embedding.
        
        Args:
            text: Text to chunk
            chunk_size: Target size of each chunk in characters
            overlap: Overlap between chunks
            respect_boundaries: Try to break at sentence/paragraph boundaries
        
        Returns:
            List of text chunks
        """
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            if end >= len(text):
                chunks.append(text[start:])
                break
            
            if respect_boundaries:
                # Try to find a good break point
                chunk = text[start:end]
                
                # Look for paragraph break
                para_break = chunk.rfind('\n\n')
                if para_break > chunk_size // 2:
                    end = start + para_break + 2
                else:
                    # Look for sentence break
                    for sep in ['. ', '! ', '? ', '.\n', '!\n', '?\n']:
                        sent_break = chunk.rfind(sep)
                        if sent_break > chunk_size // 2:
                            end = start + sent_break + len(sep)
                            break
                    else:
                        # Look for any break point
                        for sep in ['\n', ', ', '; ', ' ']:
                            break_point = chunk.rfind(sep)
                            if break_point > chunk_size // 2:
                                end = start + break_point + len(sep)
                                break
            
            chunks.append(text[start:end].strip())
            start = end - overlap
        
        return [c for c in chunks if c.strip()]
